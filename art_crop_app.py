import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from io import BytesIO

st.title("Extract Images with Titles from PowerPoint")

uploaded_file = st.file_uploader("Upload your .pptx file", type=["pptx"])

def find_best_text_below(image_shape, shapes):
    closest_text = None
    smallest_distance = float("inf")
    image_bottom = image_shape.top + image_shape.height
    image_center_x = image_shape.left + image_shape.width / 2

    for shape in shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text_top = shape.top
            text_center_x = shape.left + shape.width / 2
            vertical_gap = text_top - image_bottom
            horizontal_gap = abs(image_center_x - text_center_x)

            if 0 < vertical_gap < 1500000 and horizontal_gap < 300000:
                distance = vertical_gap + horizontal_gap
                if distance < smallest_distance:
                    smallest_distance = distance
                    closest_text = shape.text_frame.text.strip()

    return closest_text if closest_text else "Untitled"

if uploaded_file:
    prs = Presentation(uploaded_file)
    new_prs = Presentation()
    blank_slide_layout = new_prs.slide_layouts[6]
    image_slide_count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_slide_count += 1

                # Save image temporarily in memory
                image_bytes = shape.image.blob
                img_stream = BytesIO(image_bytes)

                # Get best title below image
                title_text = find_best_text_below(shape, slide.shapes)

                # Create new slide
                slide_new = new_prs.slides.add_slide(blank_slide_layout)
                pic = slide_new.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(6))
                txBox = slide_new.shapes.add_textbox(Inches(1), Inches(5.5), Inches(6), Inches(1))
                tf = txBox.text_frame
                tf.text = title_text
                tf.paragraphs[0].font.size = Pt(20)

    # Save the new presentation to memory
    output_stream = BytesIO()
    new_prs.save(output_stream)
    output_stream.seek(0)

    st.success(f"Processed {image_slide_count} images.")
    st.download_button("Download Processed PPTX", output_stream, file_name="Extracted_Images_With_Titles.pptx")
